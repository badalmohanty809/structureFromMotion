# date: 30-01-2026
# author: MohantyB
# topic: historic aerial image processing
# description: script to make ppt from historical images in a folder

import os

from osgeo import gdal
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt


# define the input his image paths
his_img_path = "D:\\datasets\\wales_gov\\penarth_head_to_cold_knap\\"
# make a list of all the paths having the image files
his_img_dir = [his_img_path + f for f in os.listdir(his_img_path)
                if '.pptx' not in f]
his_img_list = [os.listdir(f) for f in his_img_dir]

[x[0] for x in os.walk(his_img_path)]

# define output ppt path
output_ppt_path = ("D:\\datasets\\wales_gov\\penarth_head_to_cold_knap.pptx")


def make_image(in_img_path: str, out_img_path: str, scale_factor = 0.25
               , dpi =72) -> None:
    '''
    function to make low res image from high res tiff image

    in_img_path: path to input tiff image
    out_img_path: path to output low res image
    scale_factor: factor to downsample the image
    dpi: resolution of output image
    return: None
    '''
    # check the input data type to be a image file
    if not in_img_path.endswith(('.tif', '.tiff', '.img', '.png', '.jpg',
                                 '.jpeg')):
        print('Input file is not a valid image file.')
        return
    # read tiff with gdal
    dataset = gdal.Open(in_img_path)
    band = dataset.GetRasterBand(1)
    array = band.ReadAsArray()
    # downsample array to low resolution (25% of original size)
    # scale_factor = 0.25
    h, w = array.shape
    # new_h, new_w = int(h * scale_factor), int(w * scale_factor)
    array_lowres = array[::int(1/scale_factor), ::int(1/scale_factor)]
    # save the image with low resolution
    plt.imsave(out_img_path, array_lowres, cmap='gray', dpi=dpi)


def add_image_to_ppt(p: Presentation, in_img_path: str, out_img_path: str
                     )-> None:
    '''
    function to add image to ppt slide

    p: ppt presentation object
    in_img_path: path to input high res image
    out_img_path: path to input low res image
    return: None
    '''
    # add out_img_path to ppt slides, using blank slide layout
    slide = p.slides.add_slide(p.slide_layouts[6])
    # add custom title textbox
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(9)
    title_height = Inches(0.5)
    title_box = slide.shapes.add_textbox(title_left, title_top,
                title_width, title_height)
    text_frame = title_box.text_frame
    text_frame.text = in_img_path
    text_frame.paragraphs[0].font.size = Pt(14)
    #  check if the input image exists
    if not os.path.exists(out_img_path):
        print(f'Low res image file not found: {out_img_path}')
        return
    # add image below the title
    left = Inches(0.5)
    top = Inches(1.2)
    max_width = Inches(9)
    max_height = Inches(6)
    pic = slide.shapes.add_picture(out_img_path, left, top,
                width=max_width)
    # if image height exceeds max, resize by height instead
    if pic.height > max_height:
        pic.height = max_height
    #  delete the low res image file
    os.remove(out_img_path)


# Create and add title slide
p = Presentation()

# loop through all the his images and add to ppt
for i in range(len(his_img_dir)):
    print(f'Processing images in folder: {his_img_dir[i]}')
    his_img_list = [f for f in os.listdir(his_img_dir[i]) if '_lowres' not in
                    f]
    # selecting title slide layout
    s = p.slides.add_slide(p.slide_layouts[0])
    # adding title
    s.shapes.title.text = his_img_dir[i].split('\\')[-1]
    # loop through all images in the folder
    for j in range(len(his_img_list)):
        print(f'  Adding image: {his_img_list[j]}')
        in_img_path = his_img_dir[i] + '\\' + his_img_list[j]
        if os.path.isdir(in_img_path):
            his_img_list_sub = [f for f in os.listdir(in_img_path) if
                                '_lowres' not in f]
            for k in range(len(his_img_list_sub)):
                print(f'    Adding sub-image: {his_img_list_sub[k]}')
                in_img_path_sub = in_img_path + '\\' + his_img_list_sub[k]
                out_img_path_sub = (in_img_path + '\\' +
                                    his_img_list_sub[k].split('.')[0]
                                    + '_lowres.png')
                make_image(in_img_path_sub, out_img_path_sub,
                           scale_factor=0.25, dpi=70)
                add_image_to_ppt(p, in_img_path_sub, out_img_path_sub)
        else:
            out_img_path = (his_img_dir[i] + '\\' + his_img_list[j].
                            split('.')[0] + '_lowres.png')
            make_image(in_img_path, out_img_path, scale_factor=0.25, dpi=70)
            add_image_to_ppt(p, in_img_path, out_img_path)


# save the ppt
p.save(output_ppt_path)
