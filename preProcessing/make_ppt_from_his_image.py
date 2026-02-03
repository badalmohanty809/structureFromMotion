# date: 30-01-2026
# author: MohantyB
# topic: historic aerial image processing
# description: script to make ppt from historical images in a folder

import os

from osgeo import gdal
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt

# define the input his image paths
his_img_path = "D:\\datasets\\wales_gov\\penarth_head_to_cold_knap\\"
# define output ppt path
output_ppt_path = ("D:\\datasets\\wales_gov\\penarth_head_to_cold_knap.pptx")

def make_image(in_img_path: str, out_img_path: str, scale_factor = 0.25
               , dpi =72) -> None:
    '''
    function to make low res image from high res tiff image

    in_img_path: path to input tiff image
    out_img_path: path to output low res image
    scale_factor: factor to downsample the image
    dpi: resolution of output image
    return: None
    '''
    # check the input data type to be a image file
    if not in_img_path.endswith(('.tif', '.tiff', '.img', '.png', '.jpg',
                                 '.jpeg')):
        print('Input file is not a valid image file.')
        return
    # read tiff with gdal
    dataset = gdal.Open(in_img_path)
    band = dataset.GetRasterBand(1)
    array = band.ReadAsArray()
    # downsample array to low resolution (25% of original size)
    # scale_factor = 0.25
    h, w = array.shape
    # new_h, new_w = int(h * scale_factor), int(w * scale_factor)
    array_lowres = array[::int(1/scale_factor), ::int(1/scale_factor)]
    # save the image with low resolution
    plt.imsave(out_img_path, array_lowres, cmap='gray', dpi=dpi)

def add_image_to_ppt(p: Presentation, in_img_path: str, out_img_path: str
                     )-> None:
    '''
    function to add image to ppt slide

    p: ppt presentation object
    in_img_path: path to input high res image
    out_img_path: path to input low res image
    return: None
    '''
    # add out_img_path to ppt slides, using blank slide layout
    slide = p.slides.add_slide(p.slide_layouts[6])
    # add custom title textbox
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(9)
    title_height = Inches(0.5)
    title_box = slide.shapes.add_textbox(title_left, title_top,
                title_width, title_height)
    text_frame = title_box.text_frame
    text_frame.text = in_img_path
    text_frame.paragraphs[0].font.size = Pt(14)
    #  check if the input image exists
    if not os.path.exists(out_img_path):
        print(f'Low res image file not found: {out_img_path}')
        return
    # add image below the title
    left = Inches(0.5)
    top = Inches(1.2)
    max_width = Inches(9)
    max_height = Inches(6)
    pic = slide.shapes.add_picture(out_img_path, left, top,
                width=max_width)
    # if image height exceeds max, resize by height instead
    if pic.height > max_height:
        pic.height = max_height
    #  delete the low res image file
    os.remove(out_img_path)

# list out all files in the his image path
his_file_list = []
for path, subdirs, files in os.walk(his_img_path):
    for name in files:
        his_file_list.append(os.path.join(path, name))

# check for list empty condition
print(f'Number of files found: {len(his_file_list)}')
if len(his_file_list) == 0:
    print('No files found in the specified his image path.')
    exit()

# extracting the list of unique folders in the his image path
his_folder_list = list(set([f.replace(his_img_path, "").split("\\")[0]
                            for f in his_file_list]))
his_folder_list.sort()
print(f'Number of unique folders found: {len(his_folder_list)}')


# Create and add title slide
p = Presentation()
# loop through all the his images and add to ppt
for his_folder in his_folder_list:
    print(f'Found folder: {his_folder}')
    his_file_list_sub = [f for f in his_file_list if his_folder in f]
    print(f'  Number of files in folder: {len(his_file_list_sub)}')
    if len(his_file_list_sub) == 0:
        print(f'  No files found in folder: {his_folder}')
        continue
    # selecting title slide layout
    s = p.slides.add_slide(p.slide_layouts[0])
    # adding title
    s.shapes.title.text = his_folder
    # loop through all images in the folder
    for his_file in his_file_list_sub:
        print(f'  Adding image: {his_file}')
        # make low res image
        out_img_path = (his_file.split('.')[0] + '_lowres.png')
        make_image(his_file, out_img_path, scale_factor=0.25, dpi=70)
        # add image to ppt
        add_image_to_ppt(p, his_file, out_img_path)

# save the ppt
p.save(output_ppt_path)
print(f'PPT saved at: {output_ppt_path}')
